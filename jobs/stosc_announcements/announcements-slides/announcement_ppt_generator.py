import pendulum
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv
from helpers.db_helper import get_birthdays, get_anniversaries
import models.models as models

load_dotenv()


# Get birthdays and anniversaries from DB
try:
    birthdays:models.NamesList = get_birthdays()
    b_date_range=birthdays.date_range 
    birthday_data = birthdays.data
    anniversaries:models.NamesList = get_anniversaries()
    a_date_range=anniversaries.date_range
    anniversary_data = anniversaries.data
    obituary_data = []
except Exception as e:
    print(e)
    exit()


# Initialize a Presentation object based on the template
prs = Presentation("templates/announcements_template.pptx")

def print_placeholders(slide, name=""):
    for shape in slide.placeholders:
        print('%s | id=%d, name="%s"' % (name, shape.placeholder_format.idx, shape.name)) 

def get_max_names(length:int,layout:str):
    max_names:int = 10 if layout == "birthday" else 9
    variant:int = 2 if layout == "birthday" else 1
    if length <= max_names+variant:
        max_names = max_names+variant    
    if length % max_names == 1:
        max_names = max_names + 1        
    return max_names


# Function to populate birthday slides
def populate_birthday(prs, layout, date_range, names):
    max_names:int = get_max_names(len(names), "birthday")        
    for i in range(0, len(names), max_names):
        slide = prs.slides.add_slide(layout)
        # print_placeholders(slide, "birthday") 
        slide.placeholders[1].text = "\n".join(names[i:i+max_names])
        slide.placeholders[10].text = date_range

# Function to populate anniversary slides
def populate_anniversary(prs, layout, date_range, names):
    max_names:int = get_max_names(len(names), "anniversary")
    for i in range(0, len(names), max_names):
        slide = prs.slides.add_slide(layout)   
        # print_placeholders(slide, "anniversary") 
        slide.placeholders[1].text = "\n".join(names[i:i+max_names])
        slide.placeholders[10].text = date_range

# Function to populate obituary slides
def populate_obituary(prs, layout, obituaries):
    for obituary in obituaries:
        img_path = "templates/obituary_image.jpg"
        slide = prs.slides.add_slide(layout)
        # print_placeholders(slide, "obituary")
        slide.placeholders[1].insert_picture(img_path)
        slide.placeholders[2].text = obituary
        # slide.placeholders[13].text = date_range

# Get layouts
birthday_layout = prs.slide_layouts[0]  # Assuming Birthday layout is the first layout
anniversary_layout = prs.slide_layouts[1]  # Assuming Anniversary layout is the second layout
obituary_layout = prs.slide_layouts[2]  # Assuming Obituary layout is the third layout

# Populate slides
if birthday_data:
    populate_birthday(prs, birthday_layout, b_date_range, birthday_data)
if anniversary_data:
    populate_anniversary(prs, anniversary_layout, a_date_range, anniversary_data)
if obituary_data:
    populate_obituary(prs, obituary_layout, obituary_data)

# Save the presentation
file_name = f'output/Announcements'
time_string = pendulum.now("Asia/Singapore").strftime("%Y%m%d_%H%M%S")
try:
    prs.save(f"{file_name}_{time_string}.pptx")
    # TODO Change to logging
    print(f"Presentation saved as {file_name}_{time_string}.pptx")
except Exception as e:
    print(e)